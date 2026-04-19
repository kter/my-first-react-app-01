from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# Color palette
REACT_BLUE = RGBColor(0x61, 0xDA, 0xFB)   # React cyan
DARK_BG    = RGBColor(0x20, 0x23, 0x2A)   # dark background
CARD_BG    = RGBColor(0x2D, 0x31, 0x3D)   # card background
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xAA, 0xAA, 0xAA)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)
GREEN      = RGBColor(0x4E, 0xC9, 0x84)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
PURPLE     = RGBColor(0xC0, 0x82, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # completely blank

# ── helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text_box(slide, text, l, t, w, h, size=24, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def set_bg(slide, color=DARK_BG):
    add_rect(slide, 0, 0, 13.33, 7.5, color)

def accent_bar(slide, color=REACT_BLUE):
    add_rect(slide, 0, 0, 0.08, 7.5, color)

def slide_title(slide, title, subtitle=None):
    add_text_box(slide, title, 0.3, 0.25, 12.5, 0.7, size=36, bold=True, color=REACT_BLUE)
    add_rect(slide, 0.3, 1.0, 12.5, 0.04, REACT_BLUE)
    if subtitle:
        add_text_box(slide, subtitle, 0.3, 1.1, 12.5, 0.5, size=18, color=GRAY)

def code_box(slide, code, l, t, w, h, size=13):
    add_rect(slide, l, t, w, h, CARD_BG)
    txBox = slide.shapes.add_textbox(Inches(l+0.15), Inches(t+0.12),
                                      Inches(w-0.3),  Inches(h-0.24))
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split('\n'):
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = GREEN
        run.font.name = 'Courier New'

def bullet_box(slide, items, l, t, w, h, size=18, title=None, title_color=REACT_BLUE):
    if title:
        add_text_box(slide, title, l, t, w, 0.4, size=20, bold=True, color=title_color)
        t += 0.42
        h -= 0.42
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = '  • ' + item
        run.font.size = Pt(size)
        run.font.color.rgb = WHITE

# ══════════════════════════════════════════════════════════════════════════
# Slide 1 – Title
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
add_rect(s, 0, 0, 13.33, 7.5, DARK_BG)
# React logo circle
add_rect(s, 5.9, 1.2, 1.53, 1.53, REACT_BLUE)  # rough circle via square
add_text_box(s, '⚛', 5.9, 1.1, 1.6, 1.6, size=72, color=DARK_BG, align=PP_ALIGN.CENTER)

add_text_box(s, 'React 入門', 0, 3.1, 13.33, 1.0,
             size=52, bold=True, color=REACT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(s, '写経で学ぶ React の基礎から応用まで', 0, 4.2, 13.33, 0.6,
             size=24, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(s, 'Vite + React + Hooks + Jotai + SWR + Storybook', 0, 5.0, 13.33, 0.5,
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# Slide 2 – 目次
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '目次')

topics = [
    ('01', 'プロジェクト構成',   '0.5', '1.5'),
    ('02', 'State 管理',          '3.5', '1.5'),
    ('03', 'イベント処理',        '6.5', '1.5'),
    ('04', 'Hooks 基礎',          '9.5', '1.5'),
    ('05', 'Hooks 応用',          '0.5', '3.6'),
    ('06', 'スタイリング',        '3.5', '3.6'),
    ('07', '外部データ (SWR)',    '6.5', '3.6'),
    ('08', 'Jotai 状態管理',      '9.5', '3.6'),
    ('09', 'Storybook',           '0.5', '5.7'),
    ('10', 'Suspense / Lazy',     '3.5', '5.7'),
    ('11', 'パフォーマンス最適化','6.5', '5.7'),
    ('12', 'まとめ',              '9.5', '5.7'),
]

for num, label, lx, ty in topics:
    lx, ty = float(lx), float(ty)
    add_rect(s, lx, ty, 2.7, 1.5, CARD_BG)
    add_text_box(s, num, lx+0.1, ty+0.05, 0.7, 0.5, size=24, bold=True, color=REACT_BLUE)
    add_text_box(s, label, lx+0.1, ty+0.6, 2.5, 0.7, size=16, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════
# Slide 3 – プロジェクト構成
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '01. プロジェクト構成', 'Vite + React で素早く始める')

bullet_box(s, [
    'npm create vite@latest で生成',
    'エントリーポイント: src/main.jsx',
    'ルートコンポーネント: src/App.jsx',
    'Node.js 環境管理に mise を利用',
], 0.4, 1.5, 5.5, 4.0, size=19)

code_box(s, '''// main.jsx
import { StrictMode } from 'react'
import { createRoot } from 'react-dom/client'
import App from './App.jsx'

createRoot(document.getElementById('root'))
  .render(
    <StrictMode>
      <App />
    </StrictMode>
  )''', 6.2, 1.5, 6.7, 4.2)

# ══════════════════════════════════════════════════════════════════════════
# Slide 4 – State 基礎
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '02. State 管理 — useState', 'コンポーネントの内部状態を管理する')

bullet_box(s, [
    'useState(初期値) でステート変数と更新関数を取得',
    'ステートを更新すると再レンダリングが起こる',
    '親→子へ props でステートを渡せる',
    'StateTodo: 複数ステートで TODO アプリを実装',
], 0.4, 1.5, 5.8, 3.5, size=18)

code_box(s, '''// StateBasic.jsx
import { useState } from 'react';

export default function StateBasic({ init }) {
  const [count, setCount] = useState(init);
  const handleClick = () => setCount(count + 1);

  return (
    <>
      <button onClick={handleClick}>count</button>
      <p>click {count} times</p>
    </>
  );
}''', 6.2, 1.5, 6.7, 4.5)

# ══════════════════════════════════════════════════════════════════════════
# Slide 5 – StateTodo
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '02. State 管理 — Todo アプリ', '複数の State を組み合わせる')

code_box(s, '''// StateTodo.jsx（抜粋）
const [maxId, setMaxId] = useState(1);
const [title, setTitle] = useState('');
const [todo, setTodo]   = useState([]);
const [desc, setDesc]   = useState(true);

const handleClick = () => {
  setTodo([...todo, {
    id: maxId, title,
    created: new Date(), isDone: false
  }]);
  setMaxId(id => id + 1);
};

const handleDone = e => {
  setTodo(todo.map(item =>
    item.id === Number(e.target.dataset.id)
      ? { ...item, isDone: !item.isDone }
      : item
  ));
};''', 0.4, 1.5, 6.5, 5.2)

bullet_box(s, [
    '追加 / 完了 / 削除 / ソートを State だけで実装',
    '配列の更新は必ず新しい配列を返す（イミュータブル）',
    'スプレッド構文 [...todo, newItem] で追加',
    '.map() で特定要素だけ更新',
    '.filter() で要素を削除',
    'ソートは [...todo].sort() でコピーしてから',
], 7.2, 1.5, 5.8, 5.2, size=17)

# ══════════════════════════════════════════════════════════════════════════
# Slide 6 – イベント
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '03. イベント処理', 'ユーザー操作を受け取る')

bullet_box(s, [
    'onClick / onChange / onKeyDown など合成イベント',
    'ハンドラは関数参照を渡す（呼び出しではない）',
    'e.target.value でフォームの値を取得',
    'e.target.dataset.id でカスタムデータ属性を取得',
    'イベントオブジェクト (SyntheticEvent) を引数に受け取る',
], 0.4, 1.5, 5.8, 4.5, size=18)

code_box(s, '''// EventBasic
<button onClick={handleClick}>Click</button>

// EventArgs – 引数付きハンドラ
<button onClick={() => handleClick(item.id)}>

// EventKey – キーイベント
const handleKeyDown = (e) => {
  if (e.key === 'Enter') { /* ... */ }
};

// EventObj – イベントオブジェクト参照
const handleChange = (e) => {
  setTitle(e.target.value);
};''', 6.2, 1.5, 6.7, 4.5)

# ══════════════════════════════════════════════════════════════════════════
# Slide 7 – useEffect
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '04. Hooks 基礎 — useEffect', 'レンダリング後の副作用を処理する')

bullet_box(s, [
    'レンダリング後に実行される副作用処理',
    '第 2 引数の依存配列で実行タイミングを制御',
    '[] → 初回マウント時のみ',
    '[value] → value が変化した時',
    'クリーンアップ関数を return で返す（タイマー解除など）',
    'useLayoutEffect: DOM 更新前に同期実行',
], 0.4, 1.5, 5.8, 4.5, size=17)

code_box(s, '''// HookTimer.jsx
import { useEffect, useState } from 'react';

export default function HookTimer({ init }) {
  const [count, setCount] = useState(init);

  useEffect(() => {
    const t = setInterval(() => {
      setCount(c => c - 1);
    }, 1000);

    // クリーンアップ
    return () => { clearInterval(t); };
  }, []); // 初回のみ

  return <div>{count}</div>;
}''', 6.2, 1.5, 6.7, 4.8)

# ══════════════════════════════════════════════════════════════════════════
# Slide 8 – useRef
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '04. Hooks 基礎 — useRef', 'DOM 参照 & 再レンダリングを起こさない値の保持')

bullet_box(s, [
    'ref.current で値を読み書き（再レンダリングしない）',
    'DOM 要素への参照として使う',
    'タイマー ID など「状態ではない値」の保持に最適',
    'forwardRef: ref を子コンポーネントに転送',
    'useImperativeHandle: 親から子の関数を呼ぶ',
    'callback ref: ref に関数を渡す',
], 0.4, 1.5, 5.8, 4.5, size=17)

code_box(s, '''// HookRef.jsx
import { useState, useRef } from 'react';

export default function HookRef() {
  const id = useRef(null);
  const [count, setCount] = useState(0);

  const handleStart = () => {
    if (id.current === null) {
      id.current = setInterval(
        () => setCount(c => c + 1), 1000
      );
    }
  };

  const handleEnd = () => {
    clearInterval(id.current);
    id.current = null;
  };

  return (<>
    <button onClick={handleStart}>Start</button>
    <button onClick={handleEnd}>End</button>
    <p>{count}</p>
  </>);
}''', 6.2, 1.5, 6.7, 5.2)

# ══════════════════════════════════════════════════════════════════════════
# Slide 9 – useContext / useReducer
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '04. Hooks 基礎 — useContext / useReducer', 'グローバル状態とアクション駆動の状態管理')

code_box(s, '''// useContext
import MyAppContext from './MyAppContext';

export default function HookContext() {
  const config = { title: 'React', lang: 'ja-JP' };
  return (
    <MyAppContext value={config}>
      <HookContextChild />
    </MyAppContext>
  );
}''', 0.4, 1.5, 6.3, 2.5)

code_box(s, '''// useReducer
const [state, dispatch] = useReducer(
  (state, action) => {
    switch (action.type) {
      case 'update':
        return { count: state.count + 1 };
      default: return state;
    }
  },
  { count: init }
);

dispatch({ type: 'update' });''', 0.4, 4.2, 6.3, 2.8)

bullet_box(s, [
    'useContext: Props のバケツリレーを解消',
    'Provider でラップした子孫すべてが値を参照可能',
    'useReducer: 複雑な状態遷移を管理',
    'Redux ライクな action / reducer パターン',
    'dispatch(action) で状態更新を指示',
    'テーマ・言語切替など app 全体の設定に活用',
], 7.0, 1.5, 6.0, 5.5, size=17)

# ══════════════════════════════════════════════════════════════════════════
# Slide 10 – Hooks 応用 (useMemo / useCallback / useTransition)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '05. Hooks 応用 — パフォーマンス系', 'useMemo / useCallback / useTransition')

code_box(s, '''// useMemo – 計算結果をキャッシュ
const result = useMemo(() => {
  return heavyCalc(value);
}, [value]);

// useCallback – 関数参照をキャッシュ
const handleClick = useCallback(() => {
  doSomething(id);
}, [id]);

// memo – コンポーネントの再レンダリングを抑制
const Child = memo(({ value }) => {
  return <p>{value}</p>;
});''', 0.4, 1.5, 6.3, 4.0)

bullet_box(s, [
    'useMemo: 重い計算を依存値が変わるまでキャッシュ',
    'useCallback: 関数定義をメモ化（memo との組み合わせ）',
    'memo: props が変わらなければ再レンダリングしない',
    'useTransition: 優先度の低い更新を遅延',
    '  isPending フラグで UI にローディング表示',
    'useDeferredValue: 特定値の遅延バージョンを生成',
], 7.0, 1.5, 5.9, 4.5, size=17)

code_box(s, '''// useTransition
const [isPending, startTransition] = useTransition();

startTransition(() => {
  setHeavyState(newValue);
});

// useDeferredValue
const deferred = useDeferredValue(inputValue);''', 0.4, 5.7, 12.5, 1.5)

# ══════════════════════════════════════════════════════════════════════════
# Slide 11 – スタイリング
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '06. スタイリング', 'CSS から CSS-in-JS まで多様なアプローチ')

methods = [
    ('グローバル CSS',   'import ./App.css',                REACT_BLUE),
    ('CSS Modules',      'import styles from ./X.module.css', GREEN),
    ('Styled Components','styled.div`...`',                  ORANGE),
    ('Emotion',          '@emotion/react css prop',          PURPLE),
    ('Framer Motion',    'motion.div animate={...}',         YELLOW),
    ('Material UI',      'Box / Grid / Drawer コンポーネント', GRAY),
]

for i, (name, desc, color) in enumerate(methods):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.5
    ty = 1.5 + row * 1.8
    add_rect(s, lx, ty, 6.1, 1.5, CARD_BG)
    add_text_box(s, name, lx+0.15, ty+0.1, 5.8, 0.5, size=18, bold=True, color=color)
    add_text_box(s, desc, lx+0.15, ty+0.65, 5.8, 0.6, size=14, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════
# Slide 12 – SWR
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '07. 外部データ取得 — SWR', 'stale-while-revalidate キャッシュ戦略')

bullet_box(s, [
    'useSWR(key, fetcher) で宣言的にデータ取得',
    'data / isLoading / error を返す',
    '同じ key のリクエストは自動で重複排除',
    'ウィンドウフォーカス時に自動再検証',
    'Suspense / ErrorBoundary との統合も可能',
    'compare() でカスタム等値判定',
], 0.4, 1.5, 5.5, 4.5, size=18)

code_box(s, '''// SWRBasic.jsx
import useSWR from 'swr';
import fetcher from './fetcher';

export default function SWRBasic() {
  const endpoint = `https://api.openweathermap.org/...`;
  const { data, isLoading, error } =
    useSWR(endpoint, fetcher);

  if (isLoading) return <p>Loading...</p>;
  if (error) return <p>Error: {error.message}</p>;

  return (
    <figure>
      <img src={data?.weather?.[0]?.icon} />
      <figcaption>
        {data?.weather?.[0]?.description}
      </figcaption>
    </figure>
  );
}''', 6.2, 1.5, 6.7, 5.2)

# ══════════════════════════════════════════════════════════════════════════
# Slide 13 – Jotai
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '08. Jotai — アトム型状態管理', 'シンプルで柔軟なグローバル状態管理')

bullet_box(s, [
    'atom() で最小単位の状態 (atom) を定義',
    'useAtom(atom) → [value, setValue] で読み書き',
    'useResetAtom で初期値にリセット',
    'Provider でアトムの有効範囲を限定',
    'Store で複数アトムをまとめて管理',
    'localStorage 連携（atomWithStorage）',
    'Todo ロジックを atom の外に分離',
], 0.4, 1.5, 5.5, 5.0, size=17)

code_box(s, '''// atom.js
import { atom } from 'jotai';
export const counterAtom = atom(0);

// JotaiCounter.jsx
import { useAtom } from 'jotai';
import { useResetAtom } from 'jotai/utils';
import { counterAtom } from './atom';

export default function JotaiCounter() {
  const [counter, setCounter] = useAtom(counterAtom);
  const resetCounter = useResetAtom(counterAtom);

  return (<>
    <button onClick={() => setCounter(c => c + 1)}>
      Count
    </button>
    <button onClick={resetCounter}>Reset</button>
    <p>{counter} times</p>
  </>);
}''', 6.2, 1.5, 6.7, 5.2)

# ══════════════════════════════════════════════════════════════════════════
# Slide 14 – Storybook
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '09. Storybook', 'コンポーネントを独立して開発・ドキュメント化')

features = [
    ('Story の作成',     'export default / export const でストーリー定義'),
    ('argTypes',         'Controls タブで props をインタラクティブに操作'),
    ('Decorator',        'ストーリーに独自レイアウトを付与'),
    ('Viewport',         'レスポンシブ確認用のビューポート設定'),
    ('背景色設定',       'backgrounds アドオンで背景を切替'),
    ('Interaction Test', 'userEvent で操作を再現してテスト'),
    ('自動生成 Doc',     'autodocs タグで MDX ドキュメントを自動生成'),
    ('カスタム Doc',     '.mdx ファイルで独自ドキュメントを追加'),
]

for i, (title_f, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.5
    ty = 1.5 + row * 1.35
    add_rect(s, lx, ty, 6.1, 1.2, CARD_BG)
    add_text_box(s, title_f, lx+0.15, ty+0.05, 5.8, 0.45, size=16, bold=True, color=ORANGE)
    add_text_box(s, desc, lx+0.15, ty+0.55, 5.8, 0.5, size=13, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════
# Slide 15 – Suspense / Lazy
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '10. Suspense / Lazy Loading', '非同期処理と遅延読み込みの宣言的記述')

bullet_box(s, [
    'lazy(): コンポーネントを動的インポートで遅延読み込み',
    'Suspense: ローディング中の fallback を宣言',
    'ErrorBoundary: エラー時のフォールバック UI',
    'use(Promise): Promise の結果を直接コンポーネントで参照',
    'throw Promise パターン（旧 Suspense 実装）',
    'SWR + Suspense: データ取得と UI を統合',
], 0.4, 1.5, 5.8, 4.5, size=18)

code_box(s, '''// LazyBasic.jsx
import { Suspense, lazy } from 'react';

const LazyButton = lazy(() =>
  sleep(2000).then(() => import('./LazyButton'))
);

export default function LazyBasic() {
  return (
    <Suspense fallback={<p>Now Loading...</p>}>
      <LazyButton />
    </Suspense>
  );
}

// use() 関数
import { use } from 'react';
const data = use(fetchPromise);''', 6.2, 1.5, 6.7, 4.8)

# ══════════════════════════════════════════════════════════════════════════
# Slide 16 – パフォーマンスまとめ
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '11. パフォーマンス最適化まとめ', '再レンダリングを減らして UX を向上させる')

items = [
    ('memo()',          'props が変わらなければ再レンダリングしない',            REACT_BLUE),
    ('useMemo()',       '重い計算結果をキャッシュ。依存値変化時のみ再計算',       GREEN),
    ('useCallback()',   '関数オブジェクトをキャッシュ。memo と組み合わせて効果的', ORANGE),
    ('useTransition()', 'UI 更新に優先度をつけ、低優先を遅延実行',               PURPLE),
    ('useDeferredValue()','特定値の遅延バージョン。入力中のリスト更新などに有効',  YELLOW),
    ('lazy + Suspense', 'コンポーネントを必要になるまで読み込まない',             GRAY),
]

for i, (hook, desc, color) in enumerate(items):
    ty = 1.5 + i * 0.93
    add_rect(s, 0.4, ty, 12.5, 0.82, CARD_BG)
    add_text_box(s, hook, 0.6, ty+0.08, 2.8, 0.5, size=17, bold=True, color=color)
    add_text_box(s, desc, 3.5, ty+0.08, 9.0, 0.5, size=16, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════
# Slide 17 – まとめ
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s); accent_bar(s)
slide_title(s, '12. まとめ', 'この写経で学んだこと')

left_items = [
    '✅ Vite + React で高速な開発環境を構築',
    '✅ useState で状態管理の基本を習得',
    '✅ useEffect / useRef で副作用を制御',
    '✅ useContext / useReducer でグローバル状態',
    '✅ useMemo / useCallback でパフォーマンス改善',
    '✅ useTransition / useDeferredValue で優先度制御',
]
right_items = [
    '✅ CSS Modules / Emotion / Framer Motion',
    '✅ SWR で宣言的な外部データ取得',
    '✅ Jotai で軽量グローバル状態管理',
    '✅ Storybook でコンポーネント開発・文書化',
    '✅ Suspense / lazy で遅延読み込み',
    '✅ ErrorBoundary でエラーハンドリング',
]

for i, item in enumerate(left_items):
    ty = 1.5 + i * 0.7
    add_text_box(s, item, 0.4, ty, 6.2, 0.62, size=17, color=WHITE)

for i, item in enumerate(right_items):
    ty = 1.5 + i * 0.7
    add_text_box(s, item, 6.9, ty, 6.2, 0.62, size=17, color=WHITE)

add_rect(s, 0.4, 6.5, 12.5, 0.7, REACT_BLUE)
add_text_box(s, '写経を通じて React エコシステムの全体像を把握できました！', 0.4, 6.52, 12.5, 0.6,
             size=20, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────
output = '/home/user/my-first-react-app-01/output.pptx'
prs.save(output)
print(f'Saved: {output}')
