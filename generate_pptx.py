from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors (white/light theme) ──────────────────────────────
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)   # white background
C_HEADER   = RGBColor(0x00, 0x78, 0xD4)   # React-ish blue header bar
C_CODE_BG  = RGBColor(0xF3, 0xF4, 0xF6)   # light gray code block
C_DIVIDER  = RGBColor(0xCC, 0xD0, 0xD9)   # light divider
C_NOTE_BG  = RGBColor(0xFF, 0xF8, 0xE1)   # pale yellow note
C_ACCENT   = RGBColor(0x00, 0x78, 0xD4)   # blue (readable on white)
C_GREEN    = RGBColor(0x1A, 0x7F, 0x37)   # dark green
C_YELLOW   = RGBColor(0x9A, 0x6A, 0x00)   # dark amber
C_PURPLE   = RGBColor(0x6F, 0x2D, 0xA8)   # dark purple
C_RED      = RGBColor(0xC0, 0x27, 0x27)   # dark red
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)   # white (for header text)
C_GRAY     = RGBColor(0x44, 0x47, 0x54)   # dark gray body text
C_DARKROW  = RGBColor(0xF7, 0xF8, 0xFA)   # table row even
C_LIGHTROW = RGBColor(0xED, 0xEF, 0xF2)   # table row odd

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Primitive helpers ────────────────────────────────────────
def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    return slide

def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=16, bold=False, color=C_GRAY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Noto Sans JP"
    return tb

def code_block(slide, code, l, t, w, h, size=10.5):
    rect(slide, l, t, w, h, C_CODE_BG, C_DIVIDER)
    tb = slide.shapes.add_textbox(
        Inches(l+0.15), Inches(t+0.08), Inches(w-0.3), Inches(h-0.16))
    tf = tb.text_frame; tf.word_wrap = False
    first = True
    for line in code.split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = RGBColor(0x1A, 0x56, 0x76)
        r.font.name = "Courier New"

# ── Slide templates ──────────────────────────────────────────
def page_header(slide, title):
    rect(slide, 0, 0, 13.33, 0.95, C_HEADER)
    rect(slide, 0, 0, 0.38,  0.95, RGBColor(0x00, 0x50, 0x9E))
    txt(slide, title, 0.55, 0.08, 12.6, 0.8, size=26, bold=True, color=C_WHITE)

def note_bar(slide, text, t, h=0.45):
    rect(slide, 0.4, t, 12.5, h, C_NOTE_BG)
    txt(slide, f"💡 {text}", 0.55, t+0.04, 12.3, h-0.08, size=12.5, color=C_YELLOW)

# ── Reusable slide builders ──────────────────────────────────

def title_slide(title, subtitle):
    slide = new_slide()
    rect(slide, 0, 0,    13.33, 0.1,  C_ACCENT)
    rect(slide, 0, 7.4,  13.33, 0.1,  C_ACCENT)
    txt(slide, title,    0.8, 2.6, 11.7, 1.4, size=52, bold=True,
        color=C_ACCENT, align=PP_ALIGN.CENTER)
    txt(slide, subtitle, 0.8, 4.2, 11.7, 0.8, size=22,
        color=C_GRAY,   align=PP_ALIGN.CENTER)

def section_slide(num, title):
    slide = new_slide()
    rect(slide, 0, 0, 0.1, 7.5, C_ACCENT)
    rect(slide, 0.3, 3.15, 12.7, 0.07, C_DIVIDER)
    txt(slide, f"Chapter {num}", 0.5, 2.3,  12.0, 0.6, size=16, color=C_ACCENT)
    txt(slide, title,           0.5, 2.85, 12.0, 1.2, size=40, bold=True, color=RGBColor(0x1A, 0x1A, 0x2E))

def bullets_slide(title, items):
    """items: list of (text, level, style)
       style: 'h'=heading  'b'=bullet  'n'=note  's'=spacer"""
    slide = new_slide()
    page_header(slide, title)
    y = 1.05
    for text, level, style in items:
        if style == 's':
            y += 0.18; continue
        if style == 'h':
            txt(slide, text, 0.5, y, 12.5, 0.42, size=17, bold=True, color=C_GREEN)
            y += 0.45; continue
        if style == 'n':
            note_bar(slide, text, y)
            y += 0.52; continue
        # bullet
        indent = 0.4 + level * 0.45
        dot_c  = [C_ACCENT, C_GREEN, C_YELLOW][min(level, 2)]
        txt(slide, "●", indent, y+0.04, 0.3, 0.32, size=9, color=dot_c)
        fc = RGBColor(0x1A, 0x1A, 0x2E) if level == 0 else C_GRAY
        txt(slide, text, indent+0.28, y, 12.5-indent, 0.38, size=15, color=fc)
        y += 0.4

def code_slide(title, desc, code, note=None, size=10.5):
    slide = new_slide()
    page_header(slide, title)
    y = 1.0
    if desc:
        txt(slide, desc, 0.5, y, 12.5, 0.42, size=15, color=C_GRAY)
        y += 0.45
    ch = (7.5 - y - 0.1) - (0.55 if note else 0)
    code_block(slide, code, 0.4, y, 12.5, ch, size=size)
    if note:
        note_bar(slide, note, y + ch + 0.08)

def two_col_slide(title, left_items, right_items):
    """items: list of (text, bold, color_or_None)"""
    slide = new_slide()
    page_header(slide, title)
    rect(slide, 6.62, 1.05, 0.06, 6.3, C_DIVIDER)
    y = 1.1
    for text, bold, color in left_items:
        c = color or C_GRAY
        txt(slide, text, 0.4, y, 6.0, 0.42, size=14, bold=bold, color=c)
        y += 0.42
    y = 1.1
    for text, bold, color in right_items:
        c = color or C_GRAY
        txt(slide, text, 6.8, y, 6.0, 0.42, size=14, bold=bold, color=c)
        y += 0.42

def table_slide(title, headers, rows, col_widths=None):
    slide = new_slide()
    page_header(slide, title)
    n = len(headers)
    cw = col_widths or [12.5/n]*n
    lx, rh = 0.4, 0.48
    x = lx
    for h, w in zip(headers, cw):
        rect(slide, x, 1.05, w, rh, C_HEADER)
        txt(slide, h, x+0.1, 1.07, w-0.2, rh-0.06, size=13, bold=True, color=C_WHITE)
        x += w
    for ri, row in enumerate(rows):
        ry = 1.05 + rh*(ri+1)
        bg = C_DARKROW if ri%2==0 else C_LIGHTROW
        x  = lx
        for ci, (cell, w) in enumerate(zip(row, cw)):
            rect(slide, x, ry, w, rh, bg)
            fc = RGBColor(0x1A, 0x1A, 0x2E) if ci==0 else C_GRAY
            txt(slide, cell, x+0.1, ry+0.03, w-0.2, rh-0.06, size=12, color=fc)
            x += w

# ════════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════════

# ── 1. Title ─────────────────────────────────────────────────
title_slide("React 入門ガイド", "基本から実践パターンまで体系的に学ぶ")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 1  React とは
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(1, "React とは？")

bullets_slide("React の基本", [
    ("Meta（旧 Facebook）が開発した JavaScript の UI ライブラリ", 0, 'b'),
    ("「コンポーネント」という部品単位で UI を組み立てる", 0, 'b'),
    ("データが変わったとき、必要な部分だけを自動で再描画する", 0, 'b'),
    ("", 0, 's'),
    ("3 つの核心概念", 0, 'h'),
    ("コンポーネント: UI の部品。関数として書き、JSX で見た目を定義する", 0, 'b'),
    ("Props: 親から子へデータを渡す仕組み。読み取り専用", 0, 'b'),
    ("State: コンポーネント内部で管理する状態。変化すると UI が再描画される", 0, 'b'),
    ("JSX = JavaScript + HTML ライクな構文を組み合わせたもの", 0, 'n'),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 2  コンポーネントと Props
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(2, "コンポーネントと Props")

code_slide("コンポーネントの基本 / Props",
    "Props を分割代入で受け取るのが推奨パターン",
    """// Props を分割代入で受け取る（推奨パターン）
export default function MyHello({ myName }) {
  return <div>Hello, {myName}!</div>;
}

// 使う側
<MyHello myName="Taro" />""",
    "props.myName より { myName } と分割代入の方がコードが簡潔。{} の中には任意の JS 式を書ける")

code_slide("children ― 子要素を受け取る",
    "children という特別な Props でタグの中身を受け取れる",
    """export default function Panel({ children }) {
  return <div className="panel">{children}</div>;
}

// 使う側：タグの中身が children として渡される
<Panel>
  <h2>タイトル</h2>
  <p>本文テキスト</p>
</Panel>""")

code_slide("children の key による絞り込み",
    "key で子要素を識別して特定の要素だけ取り出せる（HTML の「名前付きスロット」相当）",
    """export default function TitledPanel({ children }) {
  // key が 'title' の子要素と 'body' の子要素をそれぞれ取り出す
  const title = children.find(elem => elem.key === 'title');
  const body  = children.find(elem => elem.key === 'body');

  return (
    <div className="panel">
      <div className="panel-title">{title}</div>
      <hr />
      <div className="panel-body">{body}</div>
    </div>
  );
}

// 使う側：key で役割を指定して渡す
<TitledPanel>
  <h2 key="title">パネルのタイトル</h2>
  <p  key="body">パネルの本文テキスト</p>
</TitledPanel>""",
    "子要素が 1 つのときは配列でなくオブジェクトになるため children.find() が使えない点に注意")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 3  イベント処理
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(3, "イベント処理")

code_slide("イベントハンドラの書き方",
    "onClick などのイベント属性に「関数の参照」を渡す",
    """export default function EventBasic() {
  const handleClick = () => {
    console.log(new Date().toLocaleString());
  };

  // 関数の「参照」を渡す（() なし）
  return <button onClick={handleClick}>現在時刻を表示</button>;
}

// 追加引数を渡したいとき → アロー関数でラップ
export default function EventArgs() {
  const handleClick = (e, type) => {
    const d = new Date();
    console.log(`${e.target.id}: ${type === 'date' ? d.toLocaleDateString() : d.toLocaleTimeString()}`);
  };

  return (
    <div>
      <button id="date" onClick={e => handleClick(e, 'date')}>日付</button>
      <button id="time" onClick={e => handleClick(e, 'time')}>時刻</button>
    </div>
  );
}""",
    "onClick={handleClick()} と () をつけると即座に実行される。関数の参照を渡すのが鉄則")

table_slide("主なイベントの種類",
    ["イベント", "属性", "用途"],
    [
        ["クリック",   "onClick",                    "ボタン押下"],
        ["入力変化",   "onChange",                   "テキスト入力・セレクト変更"],
        ["マウス",     "onMouseEnter / onMouseLeave", "ホバー"],
        ["キーボード", "onKeyDown / onKeyUp",         "キー入力"],
    ],
    col_widths=[2.5, 4.5, 5.5])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 4  State
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(4, "State（状態管理）")

code_slide("useState の基本",
    "コンポーネントが持つ「変化するデータ」を useState フックで宣言する",
    """import { useState } from 'react';

export default function Counter({ init }) {
  const [count, setCount] = useState(init); // 初期値は init
  return (
    <>
      <button onClick={() => setCount(count + 1)}>count</button>
      <p>{count} 回クリック</p>
    </>
  );
}""",
    "更新関数を呼ぶと React が自動で UI を再描画。let count = 0; count++ のような通常変数では更新されない")

code_slide("State のリフトアップ",
    "複数の子コンポーネントがデータを共有するとき → 共通の親に State を持たせ、更新関数を Props として渡す",
    """// 親コンポーネント
export default function StateParent() {
  const [count, setCount] = useState(0);
  const update = step => setCount(c => c + step); // 関数形式で前の値を安全に参照

  return (
    <>
      <p>合計: {count}</p>
      <Counter step={1}  onUpdate={update} />
      <Counter step={5}  onUpdate={update} />
      <Counter step={-1} onUpdate={update} />
    </>
  );
}

// 子コンポーネント：受け取った onUpdate を呼ぶだけ
function Counter({ step, onUpdate }) {
  return <button onClick={() => onUpdate(step)}>{step}</button>;
}""",
    "setCount(c => c + step) の関数形式は、非同期処理が絡む場面で前の値を確実に参照するために使う")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 5  リスト処理
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(5, "リスト処理")

code_slide("リスト処理（.map()）",
    "配列データを UI に変換するには .map() を使う",
    """import { Fragment } from 'react';

export default function BookList({ books }) {
  return (
    <dl>
      {books.map(book => (
        <Fragment key={book.isbn}>
          <dt>{book.title}（{book.price}円）</dt>
          <dd>{book.summary}</dd>
        </Fragment>
      ))}
    </dl>
  );
}

// フィルタリング：価格が 3500 円未満だけ表示
const lowPrice = books.filter(book => book.price < 3500);

// ソート：タイトル順（元の配列をコピーしてから sort）
const sorted = [...books].sort((a, b) => a.title.localeCompare(b.title));""",
    "key には配列インデックスでなくデータ固有 ID を使う。sort() は破壊的なので [...books].sort() でコピーしてから")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 6  条件付きレンダリング
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(6, "条件付きレンダリング")

code_slide("条件付きレンダリングの 3 パターン",
    "if 文 / 三項演算子 / && 演算子 を使い分ける",
    """// ① if 文 ― 最もシンプルな分岐
function Greeting({ isLoggedIn }) {
  if (isLoggedIn) return <p>ようこそ！</p>;
  return <p>ログインしてください</p>;
}

// ② 三項演算子 ― JSX の中で分岐
function Badge({ status }) {
  return (
    <span className={status === 'active' ? 'badge-green' : 'badge-gray'}>
      {status === 'active' ? '有効' : '無効'}
    </span>
  );
}

// ③ && 演算子 ― true のときだけ表示
function Notification({ hasMessage, message }) {
  return (
    <div>
      <h1>ダッシュボード</h1>
      {hasMessage && <p className="notification">{message}</p>}
    </div>
  );
}""",
    "注意: {count && <p>...</p>} で count=0 のとき「0」がそのまま表示される。{count > 0 && ...} を使う")

code_slide("コンポーネントの動的切り替え",
    "変数に入ったコンポーネント（大文字始まり）をそのままレンダリングできる",
    """import { BannerMember, BannerNew, BannerEnv } from './Banners';

export default function SelectComp() {
  const components = [BannerMember, BannerNew, BannerEnv];
  // コンポーネントを変数に代入（大文字始まりの変数名が必須）
  const SelectedComponent = components[Math.floor(Math.random() * components.length)];

  return (
    <div className="banner">
      <SelectedComponent />
    </div>
  );
}""",
    "大文字始まりの変数名が必須。小文字だと HTML タグとして解釈されてしまう")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 7  render props
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(7, "render props パターン")

code_slide("render props パターン",
    "描画内容を関数として Props で受け取る設計パターン。レイアウトと表示内容を分離できる",
    """// レイアウト担当：表示方法を render 関数として受け取る
export default function ListTemplate({ src, render }) {
  return (
    <dl>
      {src.map(elem => (
        <Fragment key={elem.isbn}>
          {render(elem)}  {/* 各要素をどう表示するかは呼び出し元が決める */}
        </Fragment>
      ))}
    </dl>
  );
}

// 使う側：同じ ListTemplate で表示内容だけ変える
<ListTemplate
  src={books}
  render={book => (
    <>
      <dt>{book.title}</dt>
      <dd>{book.price}円</dd>
    </>
  )}
/>

<ListTemplate
  src={books}
  render={book => (
    <>
      <dt>{book.title}</dt>
      <dd>著者: {book.author}</dd>
    </>
  )}
/>""",
    "children（JSX を渡す）との違い: render props は関数を渡すので、コンポーネント側のデータを呼び出し元が受け取れる")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 8  フォーム処理
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(8, "フォーム処理")

code_slide("制御されたコンポーネント",
    "入力値を必ず State で管理する「Controlled Components」パターン",
    """export default function ProfileForm() {
  const [form, setForm] = useState({ name: 'Yamada', age: 18 });

  const handleChange = e => {
    setForm({
      ...form,                          // 既存の値を保持
      [e.target.name]: e.target.value   // 変更されたフィールドだけ上書き
    });
  };

  return (
    <form>
      <input name="name" value={form.name} onChange={handleChange} />
      <input name="age"  value={form.age}  onChange={handleChange} type="number" />
      <p>Hello, {form.name}（{form.age}）</p>
    </form>
  );
}""",
    "[e.target.name] は計算プロパティ名（Computed Property Names）。フィールドが増えてもハンドラを 1 つにまとめられる")

table_slide("各入力要素のポイント",
    ["要素", "ポイント"],
    [
        ["<input type=\"text\">",      "value + onChange で制御"],
        ["<textarea>",                 "JSX では自己閉じタグ不可。value で制御"],
        ["<select>",                   "value を <select> に渡して選択状態を制御"],
        ["<input type=\"radio\">",     "checked={value === selected} で選択状態を判定"],
        ["<input type=\"checkbox\">",  "checked + onChange で制御。複数選択は配列で管理"],
        ["<input type=\"file\">",      "非制御のみ。onChange で e.target.files を取得"],
    ],
    col_widths=[3.8, 8.7])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 9  React Hook Form
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(9, "React Hook Form")

code_slide("React Hook Form の基本",
    "非制御コンポーネントベース。バリデーション付きフォームを少ないコードで実装",
    """import { useForm } from 'react-hook-form';

export default function BookForm() {
  const {
    register,
    handleSubmit,
    formState: { errors, isDirty, isValid, isSubmitting }
  } = useForm({ defaultValues: { title: '', price: '' }, mode: 'onChange' });

  const onSubmit = data => console.log(data);

  return (
    <form onSubmit={handleSubmit(onSubmit)} noValidate>
      <input
        {...register('title', {
          required: 'タイトルは必須です',
          maxLength: { value: 20, message: '20文字以内で入力してください' }
        })}
      />
      <div>{errors.title?.message}</div>

      <input
        type="email"
        {...register('email', {
          required: 'メールアドレスは必須です',
          pattern: { value: /^[^\\s@]+@[^\\s@]+\\.[^\\s@]+$/, message: '形式が正しくありません' }
        })}
      />
      <div>{errors.email?.message}</div>

      <button type="submit" disabled={!isDirty || !isValid || isSubmitting}>
        {isSubmitting ? '送信中...' : '送信'}
      </button>
    </form>
  );
}""",
    "{...register('fieldName', rules)} を展開するだけで登録完了。errors.fieldName?.message でエラー表示")

two_col_slide("制御されたコンポーネント vs React Hook Form",
    [
        ("制御されたコンポーネント", True, C_ACCENT),
        ("", False, None),
        ("管理方法", True, C_GRAY),
        ("State（useState）", False, C_WHITE),
        ("", False, None),
        ("再レンダリング", True, C_GRAY),
        ("入力のたびに発生", False, C_WHITE),
        ("", False, None),
        ("コード量", True, C_GRAY),
        ("多い", False, C_WHITE),
        ("", False, None),
        ("向いている場面", True, C_GRAY),
        ("リアルタイム連動 UI", False, C_WHITE),
    ],
    [
        ("React Hook Form", True, C_GREEN),
        ("", False, None),
        ("管理方法", True, C_GRAY),
        ("DOM 参照（非制御）", False, C_WHITE),
        ("", False, None),
        ("再レンダリング", True, C_GRAY),
        ("最小限（バリデーション時のみ）", False, C_WHITE),
        ("", False, None),
        ("コード量", True, C_GRAY),
        ("少ない", False, C_WHITE),
        ("", False, None),
        ("向いている場面", True, C_GRAY),
        ("通常の入力フォーム", False, C_WHITE),
    ])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 10  State の実践 Todo
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(10, "State の実践：Todo アプリ")

code_slide("配列 State の 3 つの操作パターン",
    "push/pop などの破壊的操作は NG。新しい配列を返す操作（map / filter / スプレッド）を使う",
    """// State 設計
// maxId ... 次の ID を管理（削除後も重複しないように連番を維持）
// title ... 入力中のタイトル
// todo  ... [{ id, title, created, isDone }, ...]
// desc  ... ソート方向のフラグ

// ① 追加
setTodo([...todo, { id: maxId, title, created: new Date(), isDone: false }]);

// ② 更新（完了トグル）
setTodo(todo.map(item =>
  item.id === targetId ? { ...item, isDone: !item.isDone } : item
));

// ③ 削除
setTodo(todo.filter(item => item.id !== targetId));""",
    "todo.push(newItem) では React が変化を検知できない。必ず新しい配列を返す操作を使う")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 11  Code Splitting & Suspense
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(11, "Code Splitting と Suspense")

code_slide("React.lazy() と <Suspense>",
    "コンポーネントを必要になったタイミングで読み込み、初期バンドルサイズを削減する",
    """import { Suspense, lazy } from 'react';

// 動的インポートで遅延読み込み
const HeavyComponent = lazy(() => import('./HeavyComponent'));

export default function App() {
  return (
    <Suspense fallback={<p>読み込み中...</p>}>
      <HeavyComponent />  {/* 読み込み完了まで fallback が表示される */}
    </Suspense>
  );
}

// 複数の lazy コンポーネントを 1 つの Suspense でまとめることも可能
const PageA = lazy(() => import('./PageA'));
const PageB = lazy(() => import('./PageB'));

<Suspense fallback={<Spinner />}>
  <PageA />
  <PageB />
</Suspense>""",
    "fallback にはローディングスピナーやスケルトン UI を渡す")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 12  use()  React 19
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(12, "use()（React 19）")

code_slide("use() で Promise を読み取る",
    "if 文や for ループの中でも呼べる。Promise が解決するまで <Suspense> が待機",
    """import { use, Suspense } from 'react';

// データを取得する Promise を作成
const fetchMessage = () =>
  fetch('/api/message').then(res => res.json()).then(data => data.text);

// use() で Promise を読み取るコンポーネント
function Message({ messagePromise }) {
  const text = use(messagePromise); // 解決するまで Suspense が待機
  return <p>{text}</p>;
}

// 呼び出し側：Suspense でラップする
export default function App() {
  return (
    <Suspense fallback={<p>Loading...</p>}>
      <Message messagePromise={fetchMessage()} />
    </Suspense>
  );
}""",
    "useEffect + fetch と異なりローディング処理を <Suspense> に任せられる。useSWR よりもシンプル")

code_slide("use() で Context を読み取る",
    "useContext の代わりに use で Context を読み取れる。条件分岐の中でも呼べるのが強み",
    """import { use, createContext } from 'react';

const ThemeContext = createContext('light');

function Button({ isSpecial }) {
  // useContext と違い、条件分岐の中で呼び出せる
  if (isSpecial) {
    const theme = use(ThemeContext);
    return <button className={theme}>特別なボタン</button>;
  }
  return <button>通常のボタン</button>;
}""",
    "use() は React 19 で追加。他のフックと異なり if 文・for ループ内でも呼び出せる")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 13  Error Boundary
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(13, "Error Boundary")

code_slide("Error Boundary の基本",
    "子コンポーネントでエラーが発生してもアプリ全体がクラッシュしないよう保護する",
    """import { ErrorBoundary } from 'react-error-boundary';

// シンプルな使い方
export default function App() {
  return (
    <ErrorBoundary fallback={<p>エラーが発生しました</p>}>
      <SomeComponent />
    </ErrorBoundary>
  );
}

// onReset でリトライボタン付きフォールバック
<ErrorBoundary
  fallback={({ resetErrorBoundary }) => (
    <div>
      <p>エラーが発生しました</p>
      <button onClick={resetErrorBoundary}>再試行</button>
    </div>
  )}
>
  <SomeComponent />
</ErrorBoundary>""",
    "注意: Error Boundary はイベントハンドラ内のエラーは捕捉しない。そちらは通常の try/catch を使う")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 14  Portal
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(14, "Portal")

code_slide("createPortal ― コンポーネントツリー外へレンダリング",
    "モーダルやツールチップなど z-index 問題を避けたいときに使う",
    """import { createPortal } from 'react-dom';

export default function Modal({ show, onClose }) {
  if (!show) return null;

  return createPortal(
    <div className="dialog">
      <p>ダイアログの内容</p>
      <button onClick={onClose}>閉じる</button>
    </div>,
    document.getElementById('dialog')  // レンダリング先の DOM 要素
  );
}

// index.html 側で Portal のマウント先を用意する
// <div id="root"></div>
// <div id="dialog"></div>  ← <body> 直下に配置""",
    "Portal を使っても、イベントのバブリングは React のコンポーネントツリーに沿って動作する")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 15  ドキュメントメタデータ（React 19）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(15, "ドキュメントメタデータ（React 19）")

code_slide("<title> / <meta> の動的操作",
    "React 19 では <title> <meta> <link> <script> をコンポーネント内に書くと自動的に <head> に配置される",
    """export default function ArticlePage({ title, description }) {
  return (
    <div>
      {/* JSX の中に書くだけで <head> に自動配置される */}
      <title>{title}</title>
      <meta name="description" content={description} />

      <h1>{title}</h1>
      <p>記事の本文...</p>
    </div>
  );
}""",
    "コンポーネントがアンマウントされると <head> から自動で取り除かれる。ページ遷移時のタイトル変更に便利")

code_slide("<script> の動的読み込み / <link> の優先度制御",
    "同じ src の <script> は重複排除。<link> の precedence で CSS 適用順を制御できる",
    """// <script> ― 同じ src を複数書いても 1 回しか読み込まれない
export function ExternalWidget() {
  return (
    <div className="widget">
      <script src="https://cdn.example.com/widget.js" async />
      <div id="widget-root" />
    </div>
  );
}

// <link> + precedence ― Suspense と連携して FOUC を防止
function StyledSection() {
  return (
    <>
      <link rel="stylesheet" href="base.css"  precedence="low" />
      <link rel="stylesheet" href="theme.css" precedence="high" />
      <div className="themed-content">コンテンツ</div>
    </>
  );
}

export default function App() {
  return (
    // CSS の読み込みが完了するまで fallback を表示
    <Suspense fallback={<p>スタイル読み込み中...</p>}>
      <StyledSection />
    </Suspense>
  );
}""",
    "precedence は任意の文字列（'low'/'high' など）で挿入順序を制御。スタイルなしコンテンツの瞬間表示（FOUC）を防げる")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 16  スタイリング手法
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(16, "スタイリング手法")

code_slide("CSS Modules / CSS-in-JS（Emotion）",
    "CSS Modules はクラス名をスコープ化。Emotion は JS の変数をそのまま CSS に使える",
    """// CSS Modules
import styles from './Panel.module.css';

export default function Panel() {
  return (
    <div className={styles.panel}>
      {/* styles.panel はビルド時に "panel_abc123" のようにハッシュ化される */}
      React is a JavaScript Library
    </div>
  );
}

// CSS-in-JS（Emotion）
/** @jsxImportSource @emotion/react */
import { css } from '@emotion/react';

export default function StyledBox({ isActive }) {
  const boxStyle = css\`
    background-color: \${isActive ? 'royalblue' : 'gray'};
    color: white;
    border-radius: 5px;
  \`;
  return <div css={boxStyle}>スタイル付きボックス</div>;
}""")

table_slide("スタイリング手法の比較",
    ["手法", "メリット", "デメリット"],
    [
        ["インラインスタイル",   "シンプル・追加パッケージ不要",     "擬似クラス（:hover など）が使えない"],
        ["CSS Modules",          "スコープが閉じる・通常の CSS",      "動的スタイルに一工夫必要"],
        ["Emotion（CSS-in-JS）", "JS の変数・条件をそのまま使える",   "バンドルサイズが増える"],
        ["Material UI (MUI)",    "完成度の高いコンポーネント群",       "カスタマイズに学習コストがある"],
    ],
    col_widths=[3.2, 4.65, 4.65])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 17  Material UI
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(17, "Material UI（MUI）")

code_slide("Material UI の基本",
    "Google の Material Design に基づいた React UI ライブラリ。ゼロから CSS を書かずに整ったデザインを実現",
    """import { Button, TextField, Box } from '@mui/material';

export default function LoginForm() {
  return (
    <Box sx={{ display: 'flex', flexDirection: 'column', gap: 2 }}>
      <TextField label="メールアドレス" variant="outlined" />
      <TextField label="パスワード" type="password" variant="outlined" />
      <Button variant="contained">ログイン</Button>
    </Box>
  );
}

// テーマのカスタマイズ・ダーク/ライトモード切り替え
import { createTheme, ThemeProvider } from '@mui/material/styles';

const darkTheme = createTheme({ palette: { mode: 'dark' } });

export default function App() {
  return (
    <ThemeProvider theme={darkTheme}>
      <LoginForm />
    </ThemeProvider>
  );
}""",
    "sx prop で inline スタイルを MUI の設計原則に沿って書ける。テーマで全体のデザインを統一可能")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 18  Framer Motion
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(18, "アニメーション（Framer Motion）")

code_slide("Framer Motion の基本",
    "宣言的にアニメーションを記述できるライブラリ",
    """import { motion, AnimatePresence } from 'framer-motion';

// 入場・ホバーアニメーション
<motion.div
  initial={{ opacity: 0, y: 20 }}
  animate={{ opacity: 1, y: 0 }}
  whileHover={{ scale: 1.05 }}
>
  アニメーション付きボックス
</motion.div>

// 削除アニメーション（AnimatePresence が必要）
<AnimatePresence>
  {show && (
    <motion.div
      initial={{ opacity: 0 }}
      animate={{ opacity: 1 }}
      exit={{ opacity: 0 }}
    >
      消えるときにフェードアウト
    </motion.div>
  )}
</AnimatePresence>

// バリアント ― 複数の状態をまとめて定義
const variants = {
  hidden:  { opacity: 0, y: 20 },
  visible: { opacity: 1, y: 0 },
};
<motion.div variants={variants} initial="hidden" animate="visible" />""",
    "AnimatePresence は要素が DOM から削除されるときの exit アニメーションに必要")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 19  useEffect
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(19, "useEffect（副作用の管理）")

code_slide("useEffect の基本",
    "レンダリング後に実行したい処理（API 呼び出し・タイマー・DOM 操作など）を登録する",
    """import { useEffect, useState } from 'react';

export default function DataFetcher({ url }) {
  const [data, setData] = useState(null);

  useEffect(() => {
    fetch(url)
      .then(res => res.json())
      .then(json => setData(json));
  }, [url]); // url が変化するたびに再実行

  return <p>{data?.title}</p>;
}

// 依存配列のパターン
// []        → 初回レンダリング後に 1 回だけ実行
// [value]   → value が変化するたびに実行
// 省略      → 毎回レンダリング後に実行（ほぼ使わない）""",
    "クリーンアップ関数（return () => clearInterval(...)）を返すと、アンマウント時に実行される")

code_slide("useLayoutEffect ― ブラウザ描画前に同期実行",
    "DOM のサイズや位置を計測してから描画したい場合に使う",
    """import { useEffect, useLayoutEffect, useState } from 'react';

export default function Example() {
  const [width, setWidth] = useState(0);

  // DOM のサイズを取得してから描画したい → useLayoutEffect
  useLayoutEffect(() => {
    setWidth(document.getElementById('box').offsetWidth);
  }, []);

  // API 呼び出しなど、画面描画を遅らせる必要がない処理 → useEffect
  useEffect(() => {
    fetch('/api/data').then(/* ... */);
  }, []);

  return <div id="box">幅: {width}px</div>;
}""",
    "useLayoutEffect は描画前に同期実行するため、SSR 環境では警告が出る。DOM 操作が必要な場面以外は useEffect を使う")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 20  SWR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(20, "SWR（データ取得ライブラリ）")

code_slide("SWR の基本",
    "Vercel 製。キャッシュ・再検証・ローディング・エラーハンドリングをシンプルな API で扱える",
    """import useSWR from 'swr';

const fetcher = url => fetch(url).then(res => res.json());

export default function WeatherWidget() {
  const { data, isLoading, error } = useSWR('/api/weather', fetcher);

  if (isLoading) return <p>Loading...</p>;
  if (error)     return <p>エラーが発生しました</p>;

  return <p>天気: {data.weather[0].description}</p>;
}

// SWRConfig ― グローバル設定
import { SWRConfig } from 'swr';

export default function App() {
  return (
    <SWRConfig value={{ fetcher, suspense: true }}>
      {/* 配下の useSWR は fetcher を自動で使う */}
      <WeatherWidget />
      <ForecastWidget />
    </SWRConfig>
  );
}""",
    "URL がキャッシュのキーになる。同じ URL を複数コンポーネントで使ってもリクエストは 1 回だけ")

code_slide("SWR + Suspense / ErrorBoundary / カスタムフック",
    "suspense: true で isLoading 判定が不要に。カスタムフックに切り出すと再利用も簡単",
    """// Suspense + ErrorBoundary との連携
export default function SWRApp() {
  return (
    <ErrorBoundary
      FallbackComponent={({ error, resetErrorBoundary }) => {
        setTimeout(resetErrorBoundary, 5000); // 5 秒後に自動リトライ
        return <p>{error.message}</p>;
      }}
    >
      <Suspense fallback={<p>Loading...</p>}>
        <SWRConfig value={{ fetcher, suspense: true }}>
          <WeatherWidget />  {/* data が undefined になることはない */}
        </SWRConfig>
      </Suspense>
    </ErrorBoundary>
  );
}

// データ取得をカスタムフックに切り出す
function useWeather() {
  const { data } = useSWR(endpoint); // 同じ URL なら複数箇所で使ってもリクエストは 1 回
  return data;
}

function WeatherIcon()    { const data = useWeather(); return <img src={...} />; }
function WeatherDetails() { const data = useWeather(); return <ul>...</ul>; }""",
    "useEffect + fetch の自前実装（ローディング・エラー・再取得の管理）を SWR は 1 行に置き換える")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 21  useRef
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(21, "useRef")

code_slide("useRef の基本",
    "値が変わっても再レンダリングしない。「レンダリングに不要だが保持したい値」に使う",
    """import { useRef, useState } from 'react';

export default function Timer() {
  const timerId = useRef(null); // setInterval の ID を保持（画面表示には不要）
  const [count, setCount] = useState(0);

  const start = () => {
    if (timerId.current === null) {
      timerId.current = setInterval(() => setCount(c => c + 1), 1000);
    }
  };

  const stop = () => {
    clearInterval(timerId.current);
    timerId.current = null;
  };

  return (
    <>
      <button onClick={start}>Start</button>
      <button onClick={stop}>Stop</button>
      <p>{count}</p>
    </>
  );
}""",
    "setInterval の ID は画面表示に使わないので State でなく Ref が適切")

code_slide("forwardRef / コールバック Ref",
    "forwardRef で親から子の DOM にアクセス。コールバック Ref でマウント時に即座に処理",
    """// forwardRef ― 親から子の DOM 要素にアクセスする
const MyTextBox = forwardRef((props, ref) => (
  <input type="text" ref={ref} {...props} />
));

export default function Parent() {
  const inputRef = useRef(null);
  return (
    <>
      <MyTextBox ref={inputRef} />
      <button onClick={() => inputRef.current.focus()}>フォーカス</button>
    </>
  );
}

// コールバック Ref ― マウント時に即座に処理（useEffect より簡潔）
export default function HookCallbackRef() {
  const [show, setShow] = useState(false);
  // ref に関数を渡す。DOM がマウントされると elem に要素が、アンマウントされると null が渡される
  const callbackRef = elem => elem?.focus();

  return (
    <>
      <button onClick={() => setShow(!show)}>展開</button>
      {show && <input ref={callbackRef} />}
    </>
  );
}""",
    "コールバック Ref は依存配列を気にする必要がない。elem?.focus() のようにオプショナルチェーンを使う")

code_slide("useImperativeHandle ― 公開するメソッドを限定",
    "forwardRef で渡すと親が子 DOM 全体にアクセスできてしまう → 公開メソッドを明示的に絞り込む",
    """import { useImperativeHandle, useRef } from 'react';

// 子コンポーネント
export default function MyTextBox({ label, ref }) {
  const input = useRef(null); // 内部の DOM 参照（親には見せない）

  useImperativeHandle(ref, () => ({
    focus() {
      input.current.focus(); // focus だけを公開。他の DOM 操作は隠蔽
    }
  }), []);

  return <label>{label}: <input type="text" ref={input} /></label>;
}

// 親コンポーネント
export default function Parent() {
  const textBoxRef = useRef(null);
  return (
    <>
      <MyTextBox label="Name" ref={textBoxRef} />
      {/* textBoxRef.current.focus() は呼べるが、DOM 直接操作はできない */}
      <button onClick={() => textBoxRef.current.focus()}>フォーカス</button>
    </>
  );
}""",
    "forwardRef が「DOM 参照をそのまま渡す」のに対し、useImperativeHandle は「特定のメソッドだけを渡す」")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 22  useCallback
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(22, "useCallback（関数のメモ化）")

code_slide("useCallback と React.memo",
    "React.memo でラップした子に関数を渡すとき、useCallback で関数をメモ化して不要な再レンダリングを防ぐ",
    """import { useState, useCallback, memo } from 'react';

// React.memo でラップした子は Props が変わらなければ再レンダリングしない
const MemoButton = memo(({ onClick, children }) => {
  console.log('MemoButton rendered'); // useCallback なしだと毎回表示される
  return <button onClick={onClick}>{children}</button>;
});

export default function Parent() {
  const [count1, setCount1] = useState(0);
  const [count2, setCount2] = useState(0);

  // count1 が変わらない限り同じ関数オブジェクトを返す
  const increment = useCallback(() => setCount1(c => c + 1), []);

  return (
    <>
      <MemoButton onClick={increment}>Count1: {count1}</MemoButton>
      <button onClick={() => setCount2(c => c + 1)}>Count2: {count2}</button>
    </>
  );
}""",
    "useCallback は React.memo とセットで使うことで効果を発揮。単体では再レンダリングを防げない")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 23  useReducer
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(23, "useReducer（複雑な State の管理）")

code_slide("useReducer の基本",
    "複数の State が絡み合う場合、State の更新ロジックを reducer にまとめられる",
    """import { useReducer } from 'react';

const reducer = (state, action) => {
  switch (action.type) {
    case 'increment': return { count: state.count + 1 };
    case 'reset':     return { count: 0 };
    default:          return state;
  }
};

export default function Counter() {
  const [state, dispatch] = useReducer(reducer, { count: 0 });

  return (
    <>
      <button onClick={() => dispatch({ type: 'increment' })}>+1</button>
      <button onClick={() => dispatch({ type: 'reset' })}>リセット</button>
      <p>{state.count}</p>
    </>
  );
}""",
    "「何が起きたか（Action）」と「どう変わるか（reducer）」を分離できるため複雑な State ロジックを整理しやすい")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 24  Context API
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(24, "Context API（グローバルな状態共有）")

code_slide("Context API の基本",
    "深くネストしたコンポーネントに途中を経由させずデータを渡せる仕組み",
    """import { createContext, useContext, useState } from 'react';

// 1. Context を作成
const ThemeContext = createContext('light');

// 2. Provider でデータを供給
export default function App() {
  const [theme, setTheme] = useState('light');
  return (
    <ThemeContext value={theme}>
      <button onClick={() => setTheme(t => t === 'light' ? 'dark' : 'light')}>
        テーマ切り替え
      </button>
      <DeepNestedComponent />
    </ThemeContext>
  );
}

// 3. どこにいても useContext で取得できる
function DeepNestedComponent() {
  const theme = useContext(ThemeContext);
  return <div className={theme}>コンテンツ</div>;
}""",
    "ダークモード・認証情報・言語設定など「アプリ全体に関わる設定」に向く。頻繁に変わるデータは Jotai などを検討")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 25  Jotai
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(25, "Jotai（グローバル状態管理ライブラリ）")

code_slide("Jotai の基本",
    "atom 単位でグローバル State を管理。useState と同じ感覚で使える",
    """// atom.js ― グローバルな State を定義
import { atom } from 'jotai';

export const countAtom = atom(0);   // グローバルなカウンター
export const todoAtom  = atom([]);  // グローバルな Todo リスト

// どのコンポーネントからでも同じ atom を共有できる
import { useAtom } from 'jotai';
import { countAtom } from './atom';

export default function Counter() {
  const [count, setCount] = useAtom(countAtom); // useState と同じ API
  return (
    <>
      <button onClick={() => setCount(c => c + 1)}>+</button>
      <p>{count}</p>
    </>
  );
}""",
    "同じ countAtom を別コンポーネントで useAtom すれば値が自動的に同期。Context と異なり atom を使うコンポーネントだけ再レンダリング")

code_slide("派生 atom・書き込み専用 atom・Provider",
    "atom は値を持つだけでなく、他の atom から計算した値や更新ロジックも持てる",
    """import { atom } from 'jotai';
import { atomWithStorage, atomWithReset, RESET } from 'jotai/utils';

// localStorage と自動同期
export const todosAtom   = atomWithStorage('todos', []);
// RESET シンボルで初期値に戻せる
export const counterAtom = atomWithReset(0);
// 読み取り専用の派生 atom
export const lastIdAtom  = atom(get => get(todosAtom).at(-1)?.id ?? 0);
// 書き込み専用 atom ― 更新ロジックをカプセル化
export const todoAddAtom = atom(null, (get, set, title) => {
  set(todosAtom, [...get(todosAtom), { id: get(lastIdAtom) + 1, title, isDone: false }]);
});

// コンポーネント側
import { useAtomValue, useSetAtom } from 'jotai';
export default function TodoList() {
  const todos   = useAtomValue(todosAtom);  // 読み取りのみ
  const todoAdd = useSetAtom(todoAddAtom);  // 書き込みのみ
  return (
    <>
      <button onClick={() => todoAdd('新しいTodo')}>追加</button>
      <ul>{todos.map(t => <li key={t.id}>{t.title}</li>)}</ul>
    </>
  );
}""",
    "useAtomValue / useSetAtom を使い分けると不要な再レンダリングをさらに抑えられる")

code_slide("Provider と createStore ― atom のスコープを分ける",
    "同じ atom でもコンポーネントツリーごとに独立した値を持たせられる",
    """import { createStore, Provider } from 'jotai';

const storeA = createStore();
const storeB = createStore();

export default function App() {
  return (
    <>
      {/* storeA の countAtom と storeB の countAtom は独立している */}
      <Provider store={storeA}>
        <Counter />  {/* storeA の値を使う */}
      </Provider>

      <Provider store={storeB}>
        <Counter />  {/* storeB の値を使う（storeA とは別） */}
      </Provider>

      {/* Provider なし：グローバルストアの値を使う */}
      <Counter />
    </>
  );
}""",
    "典型的なユースケース: 同じコンポーネントを複数置いてそれぞれ独立した状態を持つ場合。テスト環境での状態分離にも有効")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 26  Storybook
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(26, "Storybook（コンポーネントカタログ）")

code_slide("Storybook の基本",
    "コンポーネントを独立した環境でインタラクティブに確認できるツール",
    """// MyButton.stories.jsx
export default {
  component: MyButton,
  args: { label: 'ボタン' },
};

// Story を export するだけで Storybook に表示される
export const Primary = {
  args: { variant: 'primary', label: 'Primary' },
};

export const Secondary = {
  args: { variant: 'secondary', label: 'Secondary' },
};

export const Disabled = {
  args: { variant: 'primary', label: 'Disabled', disabled: true },
};

// デコレーター ― Story をラップする共通コンテキスト
export default {
  component: MyButton,
  decorators: [
    (Story) => (
      <div style={{ padding: '2rem', background: '#f0f0f0' }}>
        <Story />
      </div>
    ),
  ],
};""",
    "Props をブラウザ上で変更しながら動作確認できるため、デザイナーとの共有やコンポーネント仕様確認に役立つ")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Chapter 27  Hooks チートシート
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_slide(27, "Hooks チートシート")

table_slide("Hooks チートシート（前半）",
    ["Hook", "用途", "よくある使い所"],
    [
        ["useState",           "UI に関わる値を管理",            "カウンター・フォーム入力・モーダル開閉"],
        ["useEffect",          "レンダリング後の副作用",          "API 呼び出し・タイマー・イベントリスナー"],
        ["useLayoutEffect",    "DOM 更新後・描画前に同期実行",    "DOM サイズ計測・ちらつき防止"],
        ["useRef",             "再描画なしの値保持・DOM 参照",    "タイマー ID・フォーカス操作"],
        ["useCallback",        "関数のメモ化",                   "React.memo の子へのハンドラ渡し"],
    ],
    col_widths=[3.0, 4.0, 5.5])

table_slide("Hooks チートシート（後半）",
    ["Hook", "用途", "よくある使い所"],
    [
        ["useReducer",           "複雑な State 更新ロジック",             "State の項目が多いフォーム・Todo"],
        ["useContext",           "Context からの値取得",                  "テーマ・認証情報"],
        ["useImperativeHandle",  "子から親に公開するメソッドを限定",      "forwardRef と組み合わせて使用"],
        ["useDebugValue",        "DevTools へのデバッグ情報表示",         "カスタムフック開発時"],
        ["use（React 19）",      "Promise・Context を直接読み取る",       "Suspense と組み合わせた非同期データ取得"],
    ],
    col_widths=[3.0, 4.0, 5.5])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# まとめ
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide()
rect(slide, 0, 0,   13.33, 0.1, C_ACCENT)
rect(slide, 0, 7.4, 13.33, 0.1, C_ACCENT)
txt(slide, "まとめ：React 入門ロードマップ", 0.5, 0.25, 12.5, 0.8,
    size=30, bold=True, color=RGBColor(0x1A, 0x1A, 0x2E), align=PP_ALIGN.CENTER)

steps = [
    ("① コンポーネント・Props・JSX を理解する",       C_ACCENT),
    ("② useState でインタラクティブな UI を作る",      C_GREEN),
    ("③ useEffect で API 連携・副作用を管理する",      C_YELLOW),
    ("④ Context / Jotai でグローバル状態を管理する",   C_PURPLE),
    ("⑤ パフォーマンス最適化（memo / useCallback / lazy）", C_RED),
    ("⑥ フォーム・データ取得ライブラリを活用する",     RGBColor(0x44, 0x47, 0x54)),
]

for i, (text, color) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.25 + row * 1.65
    rect(slide, x, y, 6.1, 1.45, C_DARKROW)
    rect(slide, x, y, 0.1, 1.45, color)
    rect(slide, x, y+1.4, 6.1, 0.06, C_DIVIDER)
    txt(slide, text, x + 0.25, y + 0.42, 5.65, 0.65, size=15, bold=True, color=color)

txt(slide, "Happy Coding with React!", 0.5, 6.85, 12.5, 0.5,
    size=17, color=C_GRAY, align=PP_ALIGN.CENTER)

prs.save('output.pptx')
print(f"output.pptx を生成しました（スライド数: {len(prs.slides)} 枚）")
